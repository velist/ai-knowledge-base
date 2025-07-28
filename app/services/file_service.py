"""文件处理服务模块"""

import os
import uuid
import mimetypes
from pathlib import Path
from typing import List, Optional, Dict, Any, BinaryIO
from datetime import datetime
import asyncio
from concurrent.futures import ThreadPoolExecutor

from fastapi import UploadFile
from sqlalchemy.orm import Session
from loguru import logger
import aiofiles

# 文档处理库
import PyPDF2
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import jieba
from langdetect import detect

from app.config import settings
from app.models.database import File, KnowledgeBase, User
from app.models.schemas import FileType, FileUploadResponse, FileProcessResponse
from app.core.exceptions import FileProcessingError, ValidationError
from app.services.ai_service import AIServiceManager

class FileProcessor:
    """文件处理器基类"""
    
    def __init__(self):
        self.supported_types = []
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理文件，返回提取的内容和元数据"""
        raise NotImplementedError
    
    def can_process(self, file_type: str) -> bool:
        """检查是否支持处理该文件类型"""
        return file_type.lower() in self.supported_types

class PDFProcessor(FileProcessor):
    """PDF文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['pdf']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理PDF文件"""
        try:
            def extract_pdf_content():
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    content = ""
                    page_count = len(pdf_reader.pages)
                    
                    for page in pdf_reader.pages:
                        try:
                            content += page.extract_text() + "\n"
                        except Exception as e:
                            logger.warning(f"PDF页面提取失败: {str(e)}")
                            continue
                    
                    return {
                        "content": content.strip(),
                        "page_count": page_count,
                        "word_count": len(content.split()),
                        "metadata": {
                            "title": pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                            "author": pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                            "subject": pdf_reader.metadata.get('/Subject', '') if pdf_reader.metadata else ''
                        }
                    }
            
            # 在线程池中执行CPU密集型操作
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor() as executor:
                result = await loop.run_in_executor(executor, extract_pdf_content)
            
            return result
            
        except Exception as e:
            logger.error(f"PDF处理失败: {str(e)}")
            raise FileProcessingError(f"PDF处理失败: {str(e)}")

class DOCXProcessor(FileProcessor):
    """DOCX文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['docx']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理DOCX文件"""
        try:
            def extract_docx_content():
                doc = Document(file_path)
                
                content = ""
                for paragraph in doc.paragraphs:
                    content += paragraph.text + "\n"
                
                # 提取表格内容
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            content += cell.text + " "
                        content += "\n"
                
                return {
                    "content": content.strip(),
                    "page_count": 1,  # DOCX没有固定页数概念
                    "word_count": len(content.split()),
                    "metadata": {
                        "title": doc.core_properties.title or '',
                        "author": doc.core_properties.author or '',
                        "subject": doc.core_properties.subject or ''
                    }
                }
            
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor() as executor:
                result = await loop.run_in_executor(executor, extract_docx_content)
            
            return result
            
        except Exception as e:
            logger.error(f"DOCX处理失败: {str(e)}")
            raise FileProcessingError(f"DOCX处理失败: {str(e)}")

class PPTXProcessor(FileProcessor):
    """PPTX文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['pptx']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理PPTX文件"""
        try:
            def extract_pptx_content():
                prs = Presentation(file_path)
                
                content = ""
                slide_count = len(prs.slides)
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
                
                return {
                    "content": content.strip(),
                    "page_count": slide_count,
                    "word_count": len(content.split()),
                    "metadata": {
                        "title": prs.core_properties.title or '',
                        "author": prs.core_properties.author or '',
                        "subject": prs.core_properties.subject or ''
                    }
                }
            
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor() as executor:
                result = await loop.run_in_executor(executor, extract_pptx_content)
            
            return result
            
        except Exception as e:
            logger.error(f"PPTX处理失败: {str(e)}")
            raise FileProcessingError(f"PPTX处理失败: {str(e)}")

class XLSXProcessor(FileProcessor):
    """XLSX文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['xlsx', 'xls']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理XLSX文件"""
        try:
            def extract_xlsx_content():
                workbook = load_workbook(file_path, data_only=True)
                
                content = ""
                sheet_count = len(workbook.worksheets)
                
                for sheet in workbook.worksheets:
                    content += f"工作表: {sheet.title}\n"
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            content += row_text + "\n"
                    
                    content += "\n"
                
                return {
                    "content": content.strip(),
                    "page_count": sheet_count,
                    "word_count": len(content.split()),
                    "metadata": {
                        "title": workbook.properties.title or '',
                        "author": workbook.properties.creator or '',
                        "subject": workbook.properties.subject or ''
                    }
                }
            
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor() as executor:
                result = await loop.run_in_executor(executor, extract_xlsx_content)
            
            return result
            
        except Exception as e:
            logger.error(f"XLSX处理失败: {str(e)}")
            raise FileProcessingError(f"XLSX处理失败: {str(e)}")

class TXTProcessor(FileProcessor):
    """TXT文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['txt', 'md', 'py', 'js', 'html', 'css', 'json', 'xml']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理文本文件"""
        try:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                content = await file.read()
            
            return {
                "content": content,
                "page_count": 1,
                "word_count": len(content.split()),
                "metadata": {}
            }
            
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                async with aiofiles.open(file_path, 'r', encoding='gbk') as file:
                    content = await file.read()
                
                return {
                    "content": content,
                    "page_count": 1,
                    "word_count": len(content.split()),
                    "metadata": {}
                }
            except Exception as e:
                logger.error(f"文本文件处理失败: {str(e)}")
                raise FileProcessingError(f"文本文件处理失败: {str(e)}")
        except Exception as e:
            logger.error(f"文本文件处理失败: {str(e)}")
            raise FileProcessingError(f"文本文件处理失败: {str(e)}")

class ImageProcessor(FileProcessor):
    """图片文件处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp']
    
    async def process(self, file_path: str) -> Dict[str, Any]:
        """处理图片文件"""
        try:
            def extract_image_info():
                with Image.open(file_path) as img:
                    return {
                        "content": f"图片文件: {os.path.basename(file_path)}\n尺寸: {img.size[0]}x{img.size[1]}\n格式: {img.format}",
                        "page_count": 1,
                        "word_count": 0,
                        "metadata": {
                            "width": img.size[0],
                            "height": img.size[1],
                            "format": img.format,
                            "mode": img.mode
                        }
                    }
            
            loop = asyncio.get_event_loop()
            with ThreadPoolExecutor() as executor:
                result = await loop.run_in_executor(executor, extract_image_info)
            
            return result
            
        except Exception as e:
            logger.error(f"图片处理失败: {str(e)}")
            raise FileProcessingError(f"图片处理失败: {str(e)}")

class FileService:
    """文件服务"""
    
    def __init__(self):
        self.processors = {
            'pdf': PDFProcessor(),
            'docx': DOCXProcessor(),
            'pptx': PPTXProcessor(),
            'xlsx': XLSXProcessor(),
            'xls': XLSXProcessor(),
            'txt': TXTProcessor(),
            'md': TXTProcessor(),
            'py': TXTProcessor(),
            'js': TXTProcessor(),
            'html': TXTProcessor(),
            'css': TXTProcessor(),
            'json': TXTProcessor(),
            'xml': TXTProcessor(),
            'jpg': ImageProcessor(),
            'jpeg': ImageProcessor(),
            'png': ImageProcessor(),
            'gif': ImageProcessor(),
            'bmp': ImageProcessor(),
            'webp': ImageProcessor()
        }
        
        self.ai_service = AIServiceManager()
        
        # 确保上传目录存在
        os.makedirs(settings.upload_dir, exist_ok=True)
    
    def _get_file_type(self, filename: str) -> str:
        """获取文件类型"""
        return Path(filename).suffix.lower().lstrip('.')
    
    def _validate_file(self, file: UploadFile) -> None:
        """验证文件"""
        # 检查文件大小
        if file.size and file.size > settings.max_file_size:
            raise ValidationError(f"文件大小超过限制 ({settings.max_file_size / 1024 / 1024:.1f}MB)")
        
        # 检查文件类型
        file_type = self._get_file_type(file.filename)
        if file_type not in self.processors:
            raise ValidationError(f"不支持的文件类型: {file_type}")
    
    async def upload_file(
        self,
        db: Session,
        file: UploadFile,
        user: User,
        knowledge_base_id: int
    ) -> FileUploadResponse:
        """上传文件"""
        try:
            # 验证文件
            self._validate_file(file)
            
            # 检查知识库是否存在且用户有权限
            kb = db.query(KnowledgeBase).filter(
                KnowledgeBase.id == knowledge_base_id,
                KnowledgeBase.owner_id == user.id
            ).first()
            
            if not kb:
                raise ValidationError("知识库不存在或无权限")
            
            # 生成文件ID和路径
            file_id = str(uuid.uuid4())
            file_type = self._get_file_type(file.filename)
            file_extension = f".{file_type}" if file_type else ""
            file_path = os.path.join(settings.upload_dir, f"{file_id}{file_extension}")
            
            # 保存文件
            async with aiofiles.open(file_path, 'wb') as f:
                content = await file.read()
                await f.write(content)
            
            # 获取文件信息
            file_size = len(content)
            mime_type = mimetypes.guess_type(file.filename)[0] or 'application/octet-stream'
            
            # 创建文件记录
            db_file = File(
                file_id=file_id,
                filename=f"{file_id}{file_extension}",
                original_filename=file.filename,
                file_path=file_path,
                file_type=file_type,
                file_size=file_size,
                mime_type=mime_type,
                owner_id=user.id,
                knowledge_base_id=knowledge_base_id,
                processing_status="pending"
            )
            
            db.add(db_file)
            db.commit()
            db.refresh(db_file)
            
            # 更新知识库统计
            kb.file_count += 1
            kb.total_size += file_size
            db.commit()
            
            logger.info(f"文件上传成功: {file.filename} -> {file_id}")
            
            return FileUploadResponse(
                file_id=file_id,
                filename=file.filename,
                file_size=file_size,
                file_type=file_type,
                status="uploaded",
                message="文件上传成功"
            )
            
        except Exception as e:
            db.rollback()
            # 清理已上传的文件
            if 'file_path' in locals() and os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except:
                    pass
            
            logger.error(f"文件上传失败: {str(e)}")
            raise
    
    async def process_file(self, db: Session, file_id: str) -> FileProcessResponse:
        """处理文件"""
        try:
            # 获取文件记录
            db_file = db.query(File).filter(File.file_id == file_id).first()
            if not db_file:
                raise ValidationError("文件不存在")
            
            if db_file.is_processed:
                return FileProcessResponse(
                    file_id=file_id,
                    status="completed",
                    message="文件已处理完成",
                    content_preview=db_file.content_preview,
                    word_count=db_file.word_count,
                    page_count=db_file.page_count
                )
            
            # 更新处理状态
            db_file.processing_status = "processing"
            db.commit()
            
            # 获取处理器
            processor = self.processors.get(db_file.file_type)
            if not processor:
                raise FileProcessingError(f"不支持的文件类型: {db_file.file_type}")
            
            # 处理文件
            result = await processor.process(db_file.file_path)
            
            # 检测语言
            content = result['content']
            try:
                language = detect(content[:1000]) if content else 'unknown'
            except:
                language = 'unknown'
            
            # 更新文件记录
            db_file.is_processed = True
            db_file.processing_status = "completed"
            db_file.content_preview = content[:500] if content else ""
            db_file.page_count = result.get('page_count', 1)
            db_file.word_count = result.get('word_count', 0)
            db_file.language = language
            db.commit()
            
            # 异步生成向量嵌入
            asyncio.create_task(self._generate_embeddings(db_file.file_id, content))
            
            logger.info(f"文件处理完成: {file_id}")
            
            return FileProcessResponse(
                file_id=file_id,
                status="completed",
                message="文件处理完成",
                content_preview=db_file.content_preview,
                word_count=db_file.word_count,
                page_count=db_file.page_count,
                language=language
            )
            
        except Exception as e:
            # 更新错误状态
            if 'db_file' in locals():
                db_file.processing_status = "failed"
                db_file.processing_error = str(e)
                db.commit()
            
            logger.error(f"文件处理失败: {str(e)}")
            raise
    
    async def _generate_embeddings(self, file_id: str, content: str):
        """生成向量嵌入（异步任务）"""
        try:
            if not content.strip():
                return
            
            # 文本分块
            chunks = self._split_text(content)
            
            # 生成嵌入向量
            embeddings = []
            for chunk in chunks:
                try:
                    embedding = await self.ai_service.get_embedding(chunk)
                    embeddings.append({
                        'text': chunk,
                        'embedding': embedding,
                        'file_id': file_id
                    })
                except Exception as e:
                    logger.warning(f"生成嵌入向量失败: {str(e)}")
                    continue
            
            # TODO: 存储到向量数据库
            logger.info(f"为文件 {file_id} 生成了 {len(embeddings)} 个向量嵌入")
            
        except Exception as e:
            logger.error(f"生成向量嵌入失败: {str(e)}")
    
    def _split_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """文本分块"""
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # 如果不是最后一块，尝试在句号处分割
            if end < len(text):
                # 寻找最近的句号
                last_period = text.rfind('。', start, end)
                if last_period > start:
                    end = last_period + 1
                else:
                    # 寻找最近的空格
                    last_space = text.rfind(' ', start, end)
                    if last_space > start:
                        end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap if end < len(text) else end
        
        return chunks
    
    async def delete_file(self, db: Session, file_id: str, user: User) -> bool:
        """删除文件"""
        try:
            # 获取文件记录
            db_file = db.query(File).filter(
                File.file_id == file_id,
                File.owner_id == user.id
            ).first()
            
            if not db_file:
                raise ValidationError("文件不存在或无权限")
            
            # 删除物理文件
            if os.path.exists(db_file.file_path):
                os.remove(db_file.file_path)
            
            # 更新知识库统计
            kb = db.query(KnowledgeBase).filter(
                KnowledgeBase.id == db_file.knowledge_base_id
            ).first()
            
            if kb:
                kb.file_count = max(0, kb.file_count - 1)
                kb.total_size = max(0, kb.total_size - db_file.file_size)
            
            # 删除数据库记录
            db.delete(db_file)
            db.commit()
            
            # TODO: 删除向量数据库中的相关数据
            
            logger.info(f"文件删除成功: {file_id}")
            return True
            
        except Exception as e:
            db.rollback()
            logger.error(f"文件删除失败: {str(e)}")
            raise
    
    async def get_file_list(
        self,
        db: Session,
        user: User,
        knowledge_base_id: Optional[int] = None,
        file_type: Optional[str] = None,
        skip: int = 0,
        limit: int = 20
    ) -> List[Dict[str, Any]]:
        """获取文件列表"""
        try:
            query = db.query(File).filter(File.owner_id == user.id)
            
            if knowledge_base_id:
                query = query.filter(File.knowledge_base_id == knowledge_base_id)
            
            if file_type:
                query = query.filter(File.file_type == file_type)
            
            files = query.order_by(File.created_at.desc()).offset(skip).limit(limit).all()
            
            return [
                {
                    "file_id": file.file_id,
                    "filename": file.original_filename,
                    "file_type": file.file_type,
                    "file_size": file.file_size,
                    "is_processed": file.is_processed,
                    "processing_status": file.processing_status,
                    "word_count": file.word_count,
                    "page_count": file.page_count,
                    "language": file.language,
                    "created_at": file.created_at,
                    "updated_at": file.updated_at
                }
                for file in files
            ]
            
        except Exception as e:
            logger.error(f"获取文件列表失败: {str(e)}")
            raise

# 创建文件服务实例
file_service = FileService()

# 导出
__all__ = [
    "FileService",
    "file_service"
]